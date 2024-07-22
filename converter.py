import os
from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.util import Inches

def convert_to_pdf(image_path):
    image = Image.open(image_path)
    pdf_path = image_path.replace(".jpg", ".pdf")
    image.save(pdf_path, "PDF", resolution=100.0)
    print(f"JPG file converted to PDF: {pdf_path}")

def convert_to_word(image_path):
    document = Document()
    image = Image.open(image_path)
    width, height = image.size
    document.add_picture(image_path, width=Inches(6), height=Inches(height/width*6))
    word_path = image_path.replace(".jpg", ".docx")
    document.save(word_path)
    print(f"JPG file converted to Word: {word_path}")

def convert_to_ppt(image_path):
    prs = Presentation()
    image = Image.open(image_path)
    width, height = image.size
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    placeholder = slide.shapes.add_picture(image_path, Inches(0), Inches(0), height=Inches(6))
    ppt_path = image_path.replace(".jpg", ".pptx")
    prs.save(ppt_path)
    print(f"JPG file converted to PowerPoint: {ppt_path}")

def main():
    image_path = input("Enter the path to the JPG file: ")
    if not os.path.exists(image_path):
        print("File not found. Please check the path and try again.")
        return

    print("Select the format to convert to:")
    print("1. PDF")
    print("2. Word")
    print("3. PowerPoint")

    choice = input("Enter your choice (1-3): ")
    if choice == "1":
        convert_to_pdf(image_path)
    elif choice == "2":
        convert_to_word(image_path)
    elif choice == "3":
        convert_to_ppt(image_path)
    else:
        print("Invalid choice. Please enter 1, 2, or 3.")

if __name__ == "__main__":
    main()